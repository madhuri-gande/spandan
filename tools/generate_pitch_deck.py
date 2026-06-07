#!/usr/bin/env python3
"""Generate Spandan finalist pitch deck (Top 15) with QR codes.

Output:
  docs/Spandan_Hackathon_Pitch.pptx
  docs/assets/qr-*.png
"""
from __future__ import annotations

from pathlib import Path

import qrcode
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
DOCS = ROOT / "docs"
ASSETS = DOCS / "assets"

DEMO_URL = "http://98.84.159.117"
MAIL_URL = "http://98.84.159.117/mail"
BITLY = "https://bit.ly/4vCaAQ4"
GITHUB = "https://github.com/madhuri-gande/spandan"

RED = RGBColor(0xC8, 0x10, 0x2E)
DARK = RGBColor(0x22, 0x22, 0x2B)
GREY = RGBColor(0x5A, 0x60, 0x68)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF8, 0xF9, 0xFB)
PINK = RGBColor(0xFF, 0xDD, 0xDD)
SOFT_RED = RGBColor(0xFD, 0xE8, 0xEA)
BLUE = RGBColor(0x2A, 0x4E, 0xA1)


def _qr(path: Path, data: str) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    qrcode.make(data, box_size=10, border=2).save(path)


def _bg(slide, color: RGBColor) -> None:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def _bar(slide, top: float = 0, h: float = 0.12) -> None:
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(top), Inches(13.333), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = RED
    sh.line.fill.background()


def _title(slide, title: str, subtitle: str = "") -> None:
    tx = slide.shapes.add_textbox(Inches(0.65), Inches(0.5), Inches(12.2), Inches(1.15))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = RED
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(15)
        p2.font.color.rgb = GREY
        p2.space_before = Pt(4)


def _body(slide, paragraphs: list[str], *, top: float = 1.55, left: float = 0.7,
          width: float = 11.9, size: int = 17, color: RGBColor = DARK,
          bold_first: bool = False) -> None:
    """Full-sentence narrative blocks — not sparse bullets."""
    tx = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(5.5))
    tf = tx.text_frame
    tf.word_wrap = True
    for i, para in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = para
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(14)
        p.line_spacing = 1.25
        if bold_first and i == 0:
            p.font.bold = True


def _cols(slide, left_title: str, left_paras: list[str],
          right_title: str, right_paras: list[str], *, top: float = 1.55) -> None:
    for side, title, paras, x in [
        ("L", left_title, left_paras, 0.65),
        ("R", right_title, right_paras, 6.85),
    ]:
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(top), Inches(6.0), Inches(5.35)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = SOFT_RED
        card.line.width = Pt(1.2)
        htx = slide.shapes.add_textbox(Inches(x + 0.25), Inches(top + 0.2), Inches(5.5), Inches(0.45))
        htx.text_frame.paragraphs[0].text = title
        htx.text_frame.paragraphs[0].font.size = Pt(18)
        htx.text_frame.paragraphs[0].font.bold = True
        htx.text_frame.paragraphs[0].font.color.rgb = RED
        _body(slide, paras, top=top + 0.65, left=x + 0.25, width=5.45, size=15)


def _highlight(slide, text: str, *, top: float = 5.9) -> None:
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.65), Inches(top), Inches(12.0), Inches(0.85)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = SOFT_RED
    box.line.fill.background()
    tx = slide.shapes.add_textbox(Inches(0.85), Inches(top + 0.15), Inches(11.6), Inches(0.6))
    p = tx.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(14)
    p.font.italic = True
    p.font.color.rgb = RED
    p.alignment = PP_ALIGN.CENTER


def _stat(slide, left: float, top: float, val: str, lbl: str) -> None:
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(2.9), Inches(1.4)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = SOFT_RED
    tx = slide.shapes.add_textbox(Inches(left + 0.15), Inches(top + 0.15), Inches(2.6), Inches(1.1))
    tf = tx.text_frame
    tf.paragraphs[0].text = val
    tf.paragraphs[0].font.size = Pt(28)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RED
    p2 = tf.add_paragraph()
    p2.text = lbl
    p2.font.size = Pt(12)
    p2.font.color.rgb = GREY


def build() -> Path:
    ASSETS.mkdir(parents=True, exist_ok=True)
    qr_dash = ASSETS / "qr-dashboard.png"
    qr_mail = ASSETS / "qr-mailpit.png"
    qr_gh = ASSETS / "qr-github.png"
    qr_bitly = ASSETS / "qr-bitly.png"
    for path, url in [(qr_dash, DEMO_URL), (qr_mail, MAIL_URL), (qr_gh, GITHUB), (qr_bitly, BITLY)]:
        _qr(path, url)

    arch_svg = ASSETS / "architecture-diagram.svg"
    arch_png = ASSETS / "architecture-diagram.png"

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # ── 1 Title ──────────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    _bg(s, RED)
    _bar(s, top=6.88, h=0.62)
    box = s.shapes.add_textbox(Inches(0.7), Inches(1.35), Inches(12), Inches(4.2))
    tf = box.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for i, (txt, sz, bold, col) in enumerate([
        ("Spandan", 58, True, WHITE),
        ("Autonomous AI Blood Support Network", 26, False, WHITE),
        ("Blood Warriors Hackathon  |  Team 066  |  Top 15 Finalist", 17, False, PINK),
        ("स्पंदन — the heartbeat that keeps care in rhythm", 15, False, PINK),
    ]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = txt
        p.font.size = Pt(sz)
        p.font.bold = bold
        p.font.color.rgb = col
        p.alignment = PP_ALIGN.CENTER
        if i:
            p.space_before = Pt(12 if i == 1 else 18 if i == 2 else 8)

    # ── 2 Opening hook ───────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    _bg(s, LIGHT_BG)
    _bar(s)
    _title(s, "Imagine a coordinator at 11 PM")
    _body(s, [
        "A thalassemia patient needs a B+ transfusion in three days. "
        "The coordinator opens a spreadsheet with hundreds of donor names, "
        "starts calling — wrong language, no answer, try the next number.",
        "Meanwhile, three other patients are approaching their transfusion dates "
        "across Hyderabad. One coordinator cannot hold this rhythm alone.",
        "Spandan is the always-on layer that forecasts who needs blood, "
        "ranks the right donors, reaches them in their language, "
        "and escalates automatically — without anyone clicking refresh.",
    ])
    _highlight(s, "We built this on 7,000+ real Blood Warriors records — and it is live on AWS right now.")

    # ── 3 Problem depth ──────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    _bg(s, LIGHT_BG)
    _bar(s)
    _title(s, "The problem we are solving", "Why manual coordination breaks at scale")
    _cols(s,
          "The patient reality",
          [
              "Over 1,00,000 thalassemia patients in India need lifelong transfusions — "
              "often 500 to 700 in a lifetime, on a personal cadence every few weeks.",
              "Missing one cycle risks complications. Families live by the calendar "
              "of the next bag of blood.",
          ],
          "The coordination reality",
          [
              "Blood Warriors matches donors to patients through bridges — "
              "but outreach is still manual calls, WhatsApp forwards, and spreadsheets.",
              "Coordinators burn out. Donors get over-contacted or contacted in the wrong language. "
              "Emergencies need everyone at once; normal weeks need gentle sequential pacing.",
          ])

    # ── 4 Before / After ───────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    _bg(s, LIGHT_BG)
    _bar(s)
    _title(s, "Before Spandan  vs  After Spandan")
    _cols(s,
          "Today (manual)",
          [
              "Coordinator scans a list and guesses who to call first.",
              "Same donors called repeatedly; others never reached.",
              "No shared memory of who said yes, who timed out, who cancelled.",
              "Dashboard is a spreadsheet. Outreach stops when the laptop sleeps.",
          ],
          "With Spandan (autonomous)",
          [
              "Agent forecasts the 90-day pipeline and works most urgent patients first.",
              "ML ranks donors by likelihood to respond; blood-group rules enforce safety.",
              "Bedrock writes Telugu, Hindi, Tamil, or English emails; donors reply YES/NO via magic link.",
              "Runs 24/7 on EC2 — coordinator watches, intervenes only when needed.",
          ])

    # ── 5 Solution narrative ─────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    _bg(s, LIGHT_BG)
    _bar(s)
    _title(s, "What Spandan does", "One intelligent nervous system for the donor network")
    _body(s, [
        "Spandan is not a chatbot on a website. It is an autonomous state machine "
        "that wakes every 30 seconds, advances every urgent patient by one step, "
        "and remembers everything in DynamoDB.",
        "For each patient (bridge), it predicts when blood is needed, loads that patient's "
        "donor pool from the real dataset, ranks donors with logistic regression, "
        "and emails exactly one donor at a time — waiting for a reply before escalating.",
        "When a donor clicks YES in their inbox, the agent confirms the donation, "
        "sends a reminder, and rolls the patient's transfusion timeline forward. "
        "When they say NO or stay silent, the next-ranked donor is contacted automatically.",
        "Emergency surge mode blasts all eligible ranked donors in parallel — one shot, "
        "with guards so we never over-book a patient.",
    ])

    # ── 6 Stats ────────────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    _bg(s, LIGHT_BG)
    _bar(s)
    _title(s, "Built for real scale")
    for i, (v, l) in enumerate([
        ("7,033", "rows in organiser dataset"),
        ("6,946", "donors loaded to DynamoDB"),
        ("80", "patient bridges"),
        ("4", "outreach languages"),
        ("24/7", "agent on EC2"),
        ("~0.7", "ML ranking AUC"),
    ]):
        _stat(s, 0.65 + (i % 3) * 4.05, 1.55 + (i // 3) * 1.75, v, l)
    _body(s, [
        "Every number above comes from our deployed system — not slide fiction. "
        "Judges can open the live dashboard, watch an email arrive in MailPit, "
        "click YES, and see the KPI change in real time.",
    ], top=5.15, size=16)

    # ── 7 Forecasting ────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    _bg(s, LIGHT_BG)
    _bar(s)
    _title(s, "Demand forecasting", "Who needs blood — and when")
    _body(s, [
        "Each patient has a personal transfusion cadence stored as frequency_in_days "
        "and last_transfusion_date from the Blood Warriors dataset.",
        "Our forecasting module rolls that cadence forward: if historical dates are in the past, "
        "we project to the next future date so the dashboard shows a realistic 90-day pipeline.",
        "The autonomous agent does not work on random patients — it pulls from this urgency queue, "
        "sorted by days_until transfusion, and skips patients already served in the last 24 hours.",
        "This is rule-based clinical cadence math, not black-box ML — because transfusion timing "
        "is driven by known intervals, not guesswork.",
    ])
    _highlight(s, "Result: proactive outreach days before crisis, not reactive panic.")

    # ── 8 ML Ranking ─────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    _bg(s, LIGHT_BG)
    _bar(s)
    _title(s, "ML donor ranking", "Who to ask first — learned from real behaviour")
    _cols(s,
          "The model",
          [
              "Algorithm: Logistic Regression (scikit-learn) with StandardScaler.",
              "Label: donated_earlier from dataset — proxy for likely YES response.",
              "Trained on 80/20 split; class_weight=balanced; ROC-AUC ~0.7 on holdout.",
              "Artifact: models/donor_ranking.pkl — trained on EC2 at deploy time.",
          ],
          "Features & safety",
          [
              "Features: calls-to-donations ratio, donation history (log), eligibility, "
              "donor type (regular/one-time), call frequency.",
              "Blood compatibility is a hard rule before ML — we never score incompatible pairs.",
              "Runtime skip_score penalises donors who ignore outreach — failure learning on top of the model.",
          ])

    # ── 9 Bedrock AI ─────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    _bg(s, LIGHT_BG)
    _bar(s)
    _title(s, "Generative AI with Amazon Bedrock", "Claude Haiku 4.5 — multilingual & intent-aware")
    _body(s, [
        "Outreach: Bedrock drafts a warm, urgent-but-not-alarming donation request "
        "in the donor's preferred language — Telugu, Hindi, Tamil, or English — "
        "using patient name, hospital, and blood group as context.",
        "Reply intelligence: the same model classifies donor replies as YES, NO, CANCEL, or QUESTION — "
        "even when the text is informal or in a regional language.",
        "Follow-ups: auto-generated reminders after confirmation, polite 'already covered' messages "
        "when enough donors have said yes, and empathetic answers when donors ask questions.",
        "This is not template mail merge — every message is generated fresh, "
        "while the agent state machine decides when to send it.",
    ])

    # ── 10 Email & replies ─────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    _bg(s, LIGHT_BG)
    _bar(s)
    _title(s, "Real emails, real replies", "Not a UI mock — a full delivery pipeline")
    _body(s, [
        "When the agent decides to contact a donor, it writes the message to DynamoDB, "
        "builds an HTML email with patient details, and sends via SMTP to MailPit "
        "(demo) or Amazon SES (production) — same code path, one env var to switch.",
        "Each email contains signed magic-link buttons for YES and NO. "
        "Donors click without logging in; HMAC tokens prevent tampering.",
        "The reply lands in DynamoDB as an inbound message; the agent classifies intent "
        "and either confirms the donation or emails the next-ranked donor.",
        "Judges see the full loop in MailPit at http://98.84.159.117/mail — "
        "multilingual HTML, real SMTP, real state updates.",
    ])

    # ── 11 Agent lifecycle ───────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    _bg(s, LIGHT_BG)
    _bar(s)
    _title(s, "Agent lifecycle", "Sequential pacing with emergency override")
    steps = [
        ("Wake", "Background thread every ~30s; process pending replies first."),
        ("Select", "Pick patients from 90-day forecast queue (most urgent first)."),
        ("Rank", "Filter by blood group + eligibility; ML score top donor."),
        ("Email", "One donor per patient per wait window (default 5 min demo / 1 hr prod)."),
        ("Wait", "If no reply after window, bump skip_score and try next donor."),
        ("Confirm", "YES → donation record, reminder, advance patient cadence."),
        ("Surge", "Emergency: parallel blast to top 5 ranked donors, one-shot."),
    ]
    for i, (head, desc) in enumerate(steps):
        y = 1.48 + i * 0.78
        badge = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.65), Inches(y), Inches(1.55), Inches(0.58))
        badge.fill.solid()
        badge.fill.fore_color.rgb = RED if i % 2 == 0 else DARK
        badge.line.fill.background()
        badge.text_frame.paragraphs[0].text = head
        badge.text_frame.paragraphs[0].font.size = Pt(12)
        badge.text_frame.paragraphs[0].font.bold = True
        badge.text_frame.paragraphs[0].font.color.rgb = WHITE
        badge.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        badge.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        dtx = s.shapes.add_textbox(Inches(2.35), Inches(y + 0.05), Inches(10.3), Inches(0.55))
        dtx.text_frame.paragraphs[0].text = desc
        dtx.text_frame.paragraphs[0].font.size = Pt(15)
        dtx.text_frame.paragraphs[0].font.color.rgb = DARK

    # ── 12 Architecture image ──────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    _bg(s, WHITE)
    _bar(s)
    _title(s, "System architecture", "Live on AWS EC2 — nginx :80 → Streamlit + MailPit")
    if arch_svg.exists():
        try:
            s.shapes.add_picture(str(arch_svg), Inches(0.45), Inches(1.35), width=Inches(12.4))
        except Exception:
            _arch_text(s)
    else:
        _arch_text(s)

    # ── 13 AWS stack ─────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    _bg(s, LIGHT_BG)
    _bar(s)
    _title(s, "AWS end-to-end", "Production-shaped demo deployment")
    _cols(s,
          "Infrastructure",
          [
              "EC2 t3.small, Amazon Linux 2023, nginx on port 80.",
              "systemd service runs Streamlit + MailPit + agent 24/7.",
              "IAM role SpandanEC2Role — DynamoDB + Bedrock, no keys on disk.",
              "Security group: 80 public, 22 SSH restricted.",
          ],
          "Managed services",
          [
              "DynamoDB: donors, bridges, messages, donations, agent_log.",
              "Bedrock Runtime: Claude Haiku 4.5 in us-east-1.",
              "Optional S3 for ML model artifact; SES swap for production email.",
              "Full source on GitHub; Bitly short link for judges on mobile.",
          ])

    # ── 14 Live demo QR ────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    _bg(s, WHITE)
    _bar(s)
    _title(s, "Live demo — try it now")
    for path, left, label, url in [
        (qr_bitly, 0.55, "Share link (WhatsApp)", BITLY),
        (qr_dash, 3.55, "Dashboard", DEMO_URL),
        (qr_mail, 6.55, "MailPit inbox", MAIL_URL),
        (qr_gh, 9.55, "GitHub", GITHUB),
    ]:
        s.shapes.add_picture(str(path), Inches(left), Inches(1.65), width=Inches(2.15))
        cap = s.shapes.add_textbox(Inches(left - 0.05), Inches(3.95), Inches(2.25), Inches(0.9))
        ctf = cap.text_frame
        ctf.paragraphs[0].text = label
        ctf.paragraphs[0].font.size = Pt(12)
        ctf.paragraphs[0].font.bold = True
        ctf.paragraphs[0].font.color.rgb = RED
        ctf.paragraphs[0].alignment = PP_ALIGN.CENTER
        p2 = ctf.add_paragraph()
        p2.text = url.replace("https://", "").replace("http://", "")[:28]
        p2.font.size = Pt(8)
        p2.font.color.rgb = GREY
        p2.alignment = PP_ALIGN.CENTER

    _body(s, [
        "Login: coordinator (password shared at demo).",
        "Watch the 90-day patient pipeline, click 'Email next donor', open MailPit, "
        "click YES in the email — KPI updates without touching the keyboard again.",
    ], top=4.85, left=0.65, width=8.5, size=14)

    script = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(5.55), Inches(12.2), Inches(1.55))
    script.fill.solid()
    script.fill.fore_color.rgb = SOFT_RED
    script.line.fill.background()
    _body(s, [
        "90-second demo script: Open dashboard → show pipeline → email donor #1 → "
        "MailPit → multilingual email → YES → confirmed KPI → optional surge on one patient.",
    ], top=5.75, left=0.75, width=11.8, size=13, color=RED)

    # ── 15 Evaluation (expanded) ───────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    _bg(s, LIGHT_BG)
    _bar(s)
    _title(s, "How we map to evaluation criteria")
    criteria = [
        ("Practicality & scalability (20%)",
         "Real cadence forecasting over 90 days, per-donor cooldowns, sequential pacing to avoid "
         "over-mobilization, surge for emergencies, SES-ready email backend, DynamoDB on-demand."),
        ("Innovation (20%)",
         "Autonomous per-patient state machine — not a chatbot wrapper. Combines forecast queue, "
         "ML ranking, generative multilingual outreach, and magic-link reply loop in one agent."),
        ("AI component (20%)",
         "Bedrock Haiku for generation + intent classification; sklearn logistic regression "
         "trained on organiser dataset with 7 engineered features and ~0.7 AUC."),
        ("Real implementation (20%)",
         "Live SMTP emails, HMAC-signed replies, DynamoDB state across cycles, agent_log audit trail — "
         "judges can verify every step without us touching the UI."),
        ("End-to-end deployment (20%)",
         "GitHub repo, EC2 + nginx + systemd, IAM role, public demo URL, pitch deck with QR codes, "
         "dataset loaded to production-shaped AWS stack."),
    ]
    for i, (head, body) in enumerate(criteria):
        y = 1.42 + i * 1.12
        htx = s.shapes.add_textbox(Inches(0.65), Inches(y), Inches(12.0), Inches(1.0))
        tf = htx.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = head
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = RED
        p2 = tf.add_paragraph()
        p2.text = body
        p2.font.size = Pt(12)
        p2.font.color.rgb = DARK
        p2.space_before = Pt(3)

    # ── 16 Future & impact ─────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    _bg(s, LIGHT_BG)
    _bar(s)
    _title(s, "Impact & what comes next")
    _body(s, [
        "Near term: Blood Warriors coordinators get a single dashboard for the whole pipeline — "
        "who is due, who was contacted, who confirmed — while the agent does the repetitive outreach.",
        "Production path: swap MailPit for Amazon SES, add WhatsApp Business webhooks alongside email, "
        "scale agent workers on ECS, and connect to hospital blood-bank inventory APIs.",
        "Why it matters: every automated outreach hour returned to a coordinator is an hour "
        "they can spend with families in crisis — not re-dialing the same donor list.",
        "Spandan keeps the rhythm. Coordinators keep the humanity.",
    ], size=16)
    _highlight(s, "Live now: " + BITLY)

    # ── 17 Thank you ───────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    _bg(s, RED)
    box = s.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(8.8), Inches(3.5))
    tf = box.text_frame
    lines = [
        ("Thank you", 46, True, WHITE),
        ("Questions welcome — demo is live", 22, False, WHITE),
        (BITLY, 16, False, PINK),
        (GITHUB, 14, False, PINK),
        ("Team 066 · Blood Warriors Hackathon", 13, False, PINK),
    ]
    for i, (txt, sz, bold, col) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = txt
        p.font.size = Pt(sz)
        p.font.bold = bold
        p.font.color.rgb = col
        if i:
            p.space_before = Pt(14 if i < 3 else 10)
    s.shapes.add_picture(str(qr_bitly), Inches(10.0), Inches(1.85), width=Inches(2.35))
    stx = s.shapes.add_textbox(Inches(9.85), Inches(4.35), Inches(2.65), Inches(0.45))
    stx.text_frame.paragraphs[0].text = "Scan to open demo"
    stx.text_frame.paragraphs[0].font.size = Pt(11)
    stx.text_frame.paragraphs[0].font.color.rgb = PINK
    stx.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    out = DOCS / "Spandan_Hackathon_Pitch.pptx"
    prs.save(out)
    return out


def _arch_text(slide) -> None:
    _body(slide, [
        "Coordinator & Donor → nginx :80 → Streamlit + MailPit (/mail)",
        "Autonomous Agent: forecast → ML rank → Bedrock email → reply → DynamoDB",
        "AWS: DynamoDB (state) + Bedrock Haiku 4.5 (multilingual AI)",
        "See docs/assets/architecture-diagram.svg for full diagram.",
    ], top=1.5, size=15)


if __name__ == "__main__":
    print(f"Wrote {build()}")
